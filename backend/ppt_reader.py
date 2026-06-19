from pptx import Presentation
def extract_ppt_text(path):

    prs = Presentation(path)

    text = ""

    for slide in prs.slides:

        text += "\n===SLIDE===\n"

        slide_text = []

        for shape in slide.shapes:

            if hasattr(shape, "text_frame") and shape.text_frame:

                for para in shape.text_frame.paragraphs:

                    line = para.text.strip()

                    if line:
                        slide_text.append(line)

        text += "\n".join(slide_text)
        text += "\n"

    return text