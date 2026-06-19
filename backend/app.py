from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
import os
from pptx import Presentation
import PyPDF2

from scene_generator import generate_scenes

app = Flask(__name__)
CORS(app)

UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)


@app.route("/")
def home():
    return "Visual Learn AI Backend Running"


@app.route("/audio/<filename>")
def serve_audio(filename):
    return send_from_directory(
        "audio",
        filename,
        mimetype="audio/mpeg"
    )


# ---------- PPT READER ----------
def extract_ppt_text(path):

    prs = Presentation(path)

    text = ""

    for slide in prs.slides:

        text += "\n===SLIDE===\n"

        slide_lines = []

        # Sort shapes by vertical position (top to bottom)
        shapes = sorted(
            slide.shapes,
            key=lambda s: getattr(s, "top", 0)
        )

        for shape in shapes:

            if hasattr(shape, "text_frame") and shape.text_frame:

                for para in shape.text_frame.paragraphs:

                    line = para.text.strip()

                    if line:
                        slide_lines.append(line)

        text += "\n".join(slide_lines)
        text += "\n"

    return text


# ---------- PDF READER ----------
def extract_pdf_text(path):

    text = ""

    with open(path, "rb") as file:

        reader = PyPDF2.PdfReader(file)

        for page in reader.pages:

            page_text = page.extract_text()

            if page_text:
                text += page_text + "\n"

    return text


# ---------- TXT READER ----------
def extract_txt_text(path):

    with open(
        path,
        "r",
        encoding="utf-8"
    ) as file:

        return file.read()


@app.route("/upload", methods=["POST"])
def upload():

    if "file" not in request.files:
        return jsonify({"error": "No file received"})

    file = request.files["file"]

    if file.filename == "":
        return jsonify({"error": "No file selected"})

    filepath = os.path.join(
        UPLOAD_FOLDER,
        file.filename
    )

    file.save(filepath)

    extension = file.filename.split(".")[-1].lower()

    try:

        if extension == "pptx":
            text = extract_ppt_text(filepath)

        elif extension == "pdf":
            text = extract_pdf_text(filepath)

        elif extension == "txt":
            text = extract_txt_text(filepath)

        else:
            return jsonify({
                "error": f"Unsupported file type: {extension}"
            })

        scenes = generate_scenes(text)

        return jsonify({
            "message": "success",
            "scenes": scenes
        })

    except Exception as e:

        return jsonify({
            "error": str(e)
        })


if __name__ == "__main__":
    app.run(debug=True)